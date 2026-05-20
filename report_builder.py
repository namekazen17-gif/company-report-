"""
report_builder.py
-----------------
Builds a professional 11-slide PowerPoint report using python-pptx.
Color scheme: Deep Navy (#1B2A4A) + Electric Blue (#2D7DD2) + White + Light Grey.

Every text string is passed through sanitise_slide_text() before rendering.
Font size minimums enforced: titles 24pt, body 14pt, captions 11pt.

Slides:
  1.  Cover — company name, competitors, date
  2.  Executive Summary — who leads and why
  3.  Channel Overview — subscribers, videos, upload frequency table
  4.  Content Performance — top videos by views/engagement per company
  5.  Content Topics & Themes — what each company covers
  6.  Posting Frequency & Consistency — upload cadence analysis
  7.  Engagement Analysis — avg views, likes, comments
  8.  Gap Analysis — topics competitors cover that your company doesn't
  9.  Video Marketing Recommendations — 4 specific actions
  10. Company Score Card — ranking table with scores
  11. Summary — final takeaway
"""

import io
import re
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ---------------------------------------------------------------------------
# Text sanitiser — strips markdown before ANY text goes on a slide
# ---------------------------------------------------------------------------

def sanitise_slide_text(text, max_words=150):
    """
    Final sanitiser before text goes onto a slide.
    - Strips all markdown symbols (**, ##, *, _, backticks)
    - Removes common AI meta-phrases
    - Truncates to max_words at a sentence boundary — NEVER mid-sentence, NEVER adds '...'
    """
    if not text:
        return ""
    # Remove markdown
    text = re.sub(r'\*{1,3}', '', text)
    text = re.sub(r'#{1,6}\s?', '', text)
    text = re.sub(r'_{1,2}', '', text)
    text = re.sub(r'`{1,3}', '', text)
    text = re.sub(r'^\s*[-•]\s+', '', text, flags=re.MULTILINE)
    # Remove AI meta-phrases
    for phrase in [
        "Here is a summary:", "Main topic:", "Target audience:",
        "Key messages:", "As an AI,", "Based on the data,",
        "It is important to note that", "I have analysed"
    ]:
        text = text.replace(phrase, "")
    # Clean whitespace
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = text.strip()
    # Remove any trailing ellipsis left by AI
    text = re.sub(r'\s*\.{2,}$', '.', text)
    # Truncate to max_words at a sentence boundary — never break mid-sentence
    words = text.split()
    if len(words) > max_words:
        # Split into sentences and accumulate up to max_words
        sentences = re.split(r'(?<=[.!?])\s+', text)
        result = []
        word_count = 0
        for sentence in sentences:
            s_words = len(sentence.split())
            if word_count + s_words <= max_words:
                result.append(sentence)
                word_count += s_words
            else:
                break
        if result:
            text = ' '.join(result)
        else:
            # Single sentence longer than max_words — keep it whole, no truncation
            text = sentences[0] if sentences else text
    return text


# ---------------------------------------------------------------------------
# Minimum font size enforcement
# ---------------------------------------------------------------------------

# Slide titles: min 24pt | Body text: min 14pt | Captions/labels: min 11pt
MIN_FONT_TITLE   = 24
MIN_FONT_BODY    = 14
MIN_FONT_CAPTION = 16   # Raised: minimum 16pt for all captions/labels


def enforce_min_font(font_size, category="body"):
    """Ensures font size never drops below the minimum for its category."""
    minimums = {
        "title":   MIN_FONT_TITLE,
        "body":    MIN_FONT_BODY,
        "caption": MIN_FONT_CAPTION,
    }
    return max(font_size, minimums.get(category, MIN_FONT_CAPTION))


# ---------------------------------------------------------------------------
# Color palette
# ---------------------------------------------------------------------------

NAVY      = RGBColor(0x1B, 0x2A, 0x4A)   # Primary dark background
BLUE      = RGBColor(0x2D, 0x7D, 0xD2)   # Accent / highlight
LIGHTBLUE = RGBColor(0xD6, 0xE8, 0xF7)   # Soft fill
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREY      = RGBColor(0xF4, 0xF6, 0xF9)   # Slide background
DARKGREY  = RGBColor(0x4A, 0x4A, 0x5A)   # Body text
GOLD      = RGBColor(0xF4, 0xA2, 0x61)   # Accent callout
GREEN     = RGBColor(0x2A, 0x9D, 0x8F)   # Positive metric
RED       = RGBColor(0xE7, 0x6F, 0x51)   # Warning / gap

# ---------------------------------------------------------------------------
# Slide dimensions (16:9 widescreen)
# ---------------------------------------------------------------------------

W = Inches(13.33)
H = Inches(7.5)


# ---------------------------------------------------------------------------
# Helper functions
# ---------------------------------------------------------------------------

def fmt_num(n):
    try:
        n = int(n)
        if n >= 1_000_000:
            return f"{n/1_000_000:.1f}M"
        if n >= 1_000:
            return f"{n/1_000:.1f}K"
        return str(n)
    except Exception:
        return str(n)


def add_rect(slide, x, y, w, h, fill_color):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_textbox(slide, text, x, y, w, h,
                font_size=14, bold=False, color=DARKGREY,
                align=PP_ALIGN.LEFT, wrap=True, italic=False,
                font_category="body", max_words=None):
    """
    Adds a textbox to the slide.
    - text is sanitised before rendering
    - font_size is enforced against the minimum for its category
    - max_words overrides the sanitiser ceiling for this specific box
    """
    # Sanitise text — strip markdown, meta-phrases, truncate
    if max_words is not None:
        clean_text = sanitise_slide_text(str(text), max_words=max_words) if text else ""
    else:
        clean_text = sanitise_slide_text(str(text)) if text else ""

    # Enforce minimum font size
    safe_size = enforce_min_font(font_size, font_category)

    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf    = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text           = clean_text
    run.font.size      = Pt(safe_size)
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = color
    run.font.name      = "Calibri"
    return txBox


def add_slide_header(slide, title, subtitle=None):
    """Dark navy header bar at top of each content slide."""
    add_rect(slide, Inches(0), Inches(0), W, Inches(1.1), NAVY)
    add_textbox(slide, title, Inches(0.4), Inches(0.15), Inches(10), Inches(0.7),
                font_size=24, bold=True, color=WHITE, font_category="title")
    if subtitle:
        add_textbox(slide, subtitle, Inches(0.4), Inches(0.75), Inches(10), Inches(0.35),
                    font_size=11, color=LIGHTBLUE, font_category="caption")


def add_blue_card(slide, x, y, w, h, title, value, subtitle="", value_color=WHITE):
    """Metric card with blue background."""
    add_rect(slide, x, y, w, h, BLUE)
    add_textbox(slide, title, x + Inches(0.1), y + Inches(0.1), w - Inches(0.2), Inches(0.35),
                font_size=11, color=LIGHTBLUE, align=PP_ALIGN.CENTER, font_category="caption")
    add_textbox(slide, value, x + Inches(0.1), y + Inches(0.4), w - Inches(0.2), Inches(0.55),
                font_size=22, bold=True, color=value_color, align=PP_ALIGN.CENTER, font_category="title")
    if subtitle:
        add_textbox(slide, subtitle, x + Inches(0.1), y + Inches(0.9), w - Inches(0.2), Inches(0.3),
                    font_size=11, color=LIGHTBLUE, align=PP_ALIGN.CENTER, font_category="caption")


def add_navy_card(slide, x, y, w, h, title, body, title_color=GOLD):
    """Content card with navy background."""
    add_rect(slide, x, y, w, h, NAVY)
    add_textbox(slide, title, x + Inches(0.15), y + Inches(0.1), w - Inches(0.3), Inches(0.4),
                font_size=12, bold=True, color=title_color, font_category="caption")
    add_textbox(slide, body, x + Inches(0.15), y + Inches(0.5), w - Inches(0.3), h - Inches(0.65),
                font_size=11, color=WHITE, wrap=True, font_category="caption")


def set_slide_bg(slide, color=GREY):
    background = slide.background
    fill       = background.fill
    fill.solid()
    fill.fore_color.rgb = color


# ---------------------------------------------------------------------------
# Individual slide builders
# ---------------------------------------------------------------------------

def slide_01_cover(prs, your_company, competitors, report_date):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, Inches(0), Inches(0), W, H, NAVY)
    add_rect(slide, Inches(0), Inches(0), Inches(0.5), H, BLUE)
    add_rect(slide, Inches(0), Inches(7.1), W, Inches(0.4), GOLD)

    add_textbox(slide, "VIDEO COMPETITOR", Inches(1), Inches(1.2), Inches(10), Inches(1),
                font_size=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT, font_category="title")
    add_textbox(slide, "INTELLIGENCE REPORT", Inches(1), Inches(2.1), Inches(10), Inches(1),
                font_size=44, bold=True, color=BLUE, align=PP_ALIGN.LEFT, font_category="title")

    add_rect(slide, Inches(1), Inches(3.2), Inches(8), Inches(0.04), GOLD)

    add_textbox(slide, f"Prepared for: {your_company}", Inches(1), Inches(3.4), Inches(10), Inches(0.5),
                font_size=16, bold=True, color=WHITE, font_category="body")
    comp_text = "Competitors Analysed: " + " · ".join(competitors)
    add_textbox(slide, comp_text, Inches(1), Inches(3.95), Inches(11), Inches(0.5),
                font_size=14, color=LIGHTBLUE, font_category="body")
    add_textbox(slide, f"Report Date: {report_date}", Inches(1), Inches(4.5), Inches(8), Inches(0.4),
                font_size=14, color=GOLD, font_category="body")
    add_textbox(slide, "YouTube · AI-Powered Video Marketing Analysis",
                Inches(1), Inches(6.5), Inches(10), Inches(0.5),
                font_size=11, color=DARKGREY, italic=True, font_category="caption")


def slide_02_executive_summary(prs, your_company, insights, analysis):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_header(slide, "Executive Summary", "Top-line findings from your competitive analysis")

    leader = analysis.get("leader", "N/A")

    add_rect(slide, Inches(0.4), Inches(1.3), Inches(4.2), Inches(1.8), BLUE)
    add_textbox(slide, "MARKET LEADER", Inches(0.5), Inches(1.4), Inches(4), Inches(0.4),
                font_size=11, bold=True, color=LIGHTBLUE, align=PP_ALIGN.CENTER, font_category="caption")
    add_textbox(slide, leader, Inches(0.5), Inches(1.8), Inches(4), Inches(0.7),
                font_size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_category="title")
    add_textbox(slide, "in YouTube Video Marketing", Inches(0.5), Inches(2.45), Inches(4), Inches(0.4),
                font_size=11, color=LIGHTBLUE, align=PP_ALIGN.CENTER, font_category="caption")

    summary = insights.get("executive_summary", "")
    add_rect(slide, Inches(4.9), Inches(1.3), Inches(8), Inches(1.8), WHITE)
    add_textbox(slide, summary, Inches(5.1), Inches(1.4), Inches(7.6), Inches(1.6),
                font_size=14, color=DARKGREY, wrap=True, font_category="body")

    ranked   = analysis.get("ranked", [])
    card_w   = Inches(2.4)
    card_gap = Inches(0.25)
    start_x  = Inches(0.4)
    y        = Inches(3.4)
    for i, r in enumerate(ranked[:4]):
        x          = start_x + i * (card_w + card_gap)
        rank_label = ["🥇 #1", "🥈 #2", "🥉 #3", "#4"][i] if i < 4 else f"#{i+1}"
        add_blue_card(slide, x, y, card_w, Inches(1.5),
                      r["company"], f"{r['total_score']}/10", rank_label)

    leader_text = insights.get("leader_analysis", "")
    add_rect(slide, Inches(0.4), Inches(5.1), Inches(12.5), Inches(2.0), WHITE)
    add_textbox(slide, "Why The Leader Wins", Inches(0.6), Inches(5.2), Inches(12), Inches(0.4),
                font_size=14, bold=True, color=NAVY, font_category="body")
    add_textbox(slide, leader_text, Inches(0.6), Inches(5.65), Inches(12), Inches(1.3),
                font_size=14, color=DARKGREY, wrap=True, font_category="body")


def slide_03_channel_overview(prs, all_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_header(slide, "Channel Overview Comparison",
                     "Subscribers · Total Videos · Total Views · Upload Frequency")

    cols       = ["Company", "Subscribers", "Total Videos", "Total Views", "Avg Views/Video", "Upload Freq"]
    col_widths = [Inches(2.8), Inches(1.7), Inches(1.7), Inches(1.8), Inches(1.9), Inches(1.7)]
    start_x    = Inches(0.35)
    row_h      = Inches(0.55)
    header_y   = Inches(1.3)

    x = start_x
    for col, cw in zip(cols, col_widths):
        add_rect(slide, x, header_y, cw - Inches(0.05), row_h, NAVY)
        add_textbox(slide, col, x + Inches(0.05), header_y + Inches(0.1),
                    cw - Inches(0.1), Inches(0.4),
                    font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                    font_category="caption")
        x += cw

    found_data = [d for d in all_data if d["channel_found"]]
    for row_i, d in enumerate(found_data):
        y  = header_y + row_h + row_i * row_h
        bg = GREY if row_i % 2 == 0 else WHITE
        x  = start_x
        freq      = d.get("upload_frequency_days")
        freq_text = f"Every {freq}d" if freq else "N/A"
        values = [
            d["company_name"],
            fmt_num(d["subscribers"]),
            fmt_num(d["total_videos"]),
            fmt_num(d["total_views"]),
            fmt_num(d["avg_views_per_video"]),
            freq_text
        ]
        for i, (val, cw) in enumerate(zip(values, col_widths)):
            add_rect(slide, x, y, cw - Inches(0.05), row_h, bg)
            fc = BLUE if i == 0 else DARKGREY
            fb = i == 0
            add_textbox(slide, val, x + Inches(0.05), y + Inches(0.1),
                        cw - Inches(0.1), Inches(0.4),
                        font_size=11, bold=fb, color=fc, align=PP_ALIGN.CENTER,
                        font_category="caption")
            x += cw

    not_found = [d for d in all_data if not d["channel_found"]]
    if not_found:
        note = "No YouTube channel found for: " + ", ".join(d["company_name"] for d in not_found)
        add_textbox(slide, note, Inches(0.35), Inches(6.9), Inches(12), Inches(0.4),
                    font_size=11, color=RED, italic=True, font_category="caption")


def slide_04_content_performance(prs, all_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_header(slide, "Content Performance",
                     "Top performing videos by views and engagement")

    found   = [d for d in all_data if d["channel_found"] and d.get("top_videos")]
    card_w  = Inches(12.5) / max(len(found), 1)
    start_x = Inches(0.35)

    for i, d in enumerate(found):
        x = start_x + i * card_w
        add_rect(slide, x, Inches(1.25), card_w - Inches(0.1), Inches(0.45), NAVY)
        add_textbox(slide, d["company_name"], x + Inches(0.05), Inches(1.3),
                    card_w - Inches(0.15), Inches(0.35),
                    font_size=12, bold=True, color=WHITE, font_category="body")
        y = Inches(1.75)
        for j, v in enumerate(d["top_videos"][:4]):
            title_txt = v["title"][:60] + ("…" if len(v["title"]) > 60 else "")
            # White card background
            add_rect(slide, x, y, card_w - Inches(0.1), Inches(1.22), WHITE)
            # Blue left accent bar
            add_rect(slide, x, y, Inches(0.04), Inches(1.22), BLUE)
            # Video title
            add_textbox(slide, title_txt,
                        x + Inches(0.1), y + Inches(0.06), card_w - Inches(0.22), Inches(0.5),
                        font_size=11, bold=True, color=NAVY, wrap=True, font_category="body")
            # Metrics line — pushed lower so it never overlaps the title
            metrics = f"👁 {fmt_num(v['views'])}   ❤ {fmt_num(v['likes'])}   💬 {fmt_num(v['comments'])}"
            add_textbox(slide, metrics,
                        x + Inches(0.1), y + Inches(0.82), card_w - Inches(0.22), Inches(0.28),
                        font_size=10, color=DARKGREY, font_category="body")
            y += Inches(1.25)


def slide_05_content_topics(prs, all_data, analysis):
    """Slide 5: Topic list — one topic per row per company, no word breaking."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_header(slide, "Content Topics & Themes",
                     "Key topics each company covers — derived from video content analysis")

    found   = [d for d in all_data if d["channel_found"]]
    card_w  = Inches(12.5) / max(len(found), 1)
    start_x = Inches(0.35)

    for i, d in enumerate(found):
        x      = start_x + i * card_w
        topics = d.get("topics", [])

        # Navy card background
        add_rect(slide, x, Inches(1.25), card_w - Inches(0.1), Inches(5.9), NAVY)

        # Company name title
        add_textbox(slide, d["company_name"],
                    x + Inches(0.15), Inches(1.35),
                    card_w - Inches(0.3), Inches(0.45),
                    font_size=14, bold=True, color=GOLD, font_category="body")

        # Topics as a clean vertical list — one per row, full card width, no truncation
        row_h   = Inches(0.42)
        row_gap = Inches(0.06)
        row_y   = Inches(1.9)

        for j, topic in enumerate(topics[:10]):
            if row_y + row_h > Inches(6.9):
                break
            # Alternating subtle background
            row_bg = BLUE if j % 2 == 0 else RGBColor(0x1A, 0x4A, 0x9A)
            add_rect(slide, x + Inches(0.12), row_y, card_w - Inches(0.3), row_h, row_bg)
            add_textbox(slide, topic,
                        x + Inches(0.22), row_y + Inches(0.05),
                        card_w - Inches(0.5), row_h - Inches(0.08),
                        font_size=12, bold=False, color=WHITE,
                        align=PP_ALIGN.LEFT, wrap=False, font_category="body")
            row_y += row_h + row_gap


def _slide_content_insight_your_company(prs, your_company_data, insight_text):
    """
    Insight slide A: Full-width horizontal insight for YOUR company.
    Large hero layout — company name prominent, insight text big and readable.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_header(slide, "Content Strategy Insight",
                     f"{your_company_data.get('company_name', '')} — AI-analysed content intelligence")

    d = your_company_data

    # Full-width navy background panel
    add_rect(slide, Inches(0.35), Inches(1.3), Inches(12.6), Inches(5.85), NAVY)

    # Gold accent bar left side
    add_rect(slide, Inches(0.35), Inches(1.3), Inches(0.08), Inches(5.85), GOLD)

    # YOUR COMPANY label
    add_textbox(slide, "YOUR COMPANY", Inches(0.65), Inches(1.5), Inches(6), Inches(0.35),
                font_size=11, bold=True, color=GOLD, font_category="caption")

    # Company name — large
    add_textbox(slide, d.get("company_name", ""), Inches(0.65), Inches(1.85), Inches(8), Inches(0.7),
                font_size=28, bold=True, color=WHITE, font_category="title")

    # Stats row
    subs  = fmt_num(d.get("subscribers", 0))
    vids  = fmt_num(d.get("total_videos", 0))
    views = fmt_num(d.get("avg_views_per_video", 0))
    freq  = d.get("upload_frequency_days")
    freq_txt = f"Every {freq}d" if freq else "N/A"

    stats = [
        ("Subscribers", subs),
        ("Total Videos", vids),
        ("Avg Views/Video", views),
        ("Upload Freq", freq_txt),
    ]
    stat_w = Inches(2.8)
    for si, (label, val) in enumerate(stats):
        sx = Inches(0.65) + si * stat_w
        add_rect(slide, sx, Inches(2.65), stat_w - Inches(0.1), Inches(0.75), BLUE)
        add_textbox(slide, label, sx + Inches(0.08), Inches(2.68),
                    stat_w - Inches(0.18), Inches(0.25),
                    font_size=10, color=LIGHTBLUE, align=PP_ALIGN.CENTER, font_category="caption")
        add_textbox(slide, val, sx + Inches(0.08), Inches(2.9),
                    stat_w - Inches(0.18), Inches(0.4),
                    font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_category="body")

    # Divider
    add_rect(slide, Inches(0.65), Inches(3.55), Inches(12.0), Inches(0.03), GOLD)

    # CONTENT INSIGHT label
    add_textbox(slide, "CONTENT STRATEGY INSIGHT", Inches(0.65), Inches(3.65),
                Inches(6), Inches(0.3),
                font_size=11, bold=True, color=GOLD, font_category="caption")

    # Insight text — large, full width, readable
    add_textbox(slide, insight_text, Inches(0.65), Inches(4.0), Inches(12.0), Inches(2.9),
                font_size=15, color=WHITE, wrap=True, font_category="body", max_words=130)


def _slide_content_insight_pair(prs, company_a, insight_a, company_b, insight_b):
    """
    Insight slide B/C: Two companies side by side, each with full insight.
    company_b / insight_b can be None if there's an odd number of companies.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_header(slide, "Content Strategy Insights",
                     "AI-analysed content intelligence — competitor deep dive")

    half_w = Inches(6.15)

    for col, (comp_data, insight) in enumerate([
        (company_a, insight_a),
        (company_b, insight_b)
    ]):
        if comp_data is None:
            continue
        ox = Inches(0.35) + col * (half_w + Inches(0.1))

        # Card background
        add_rect(slide, ox, Inches(1.3), half_w, Inches(5.85), NAVY)

        # Blue accent top bar
        add_rect(slide, ox, Inches(1.3), half_w, Inches(0.06), BLUE)

        # Company name
        add_textbox(slide, comp_data.get("company_name", ""),
                    ox + Inches(0.15), Inches(1.45),
                    half_w - Inches(0.3), Inches(0.5),
                    font_size=18, bold=True, color=GOLD, font_category="body")

        # Quick stats
        subs     = fmt_num(comp_data.get("subscribers", 0))
        avg_v    = fmt_num(comp_data.get("avg_views_per_video", 0))
        eng      = comp_data.get("engagement_rate", 0)
        freq     = comp_data.get("upload_frequency_days")
        freq_txt = f"Every {freq}d" if freq else "N/A"

        stats_text = f"Subscribers: {subs}  |  Avg Views: {avg_v}  |  Engagement: {eng}%  |  Freq: {freq_txt}"
        add_rect(slide, ox + Inches(0.15), Inches(2.05), half_w - Inches(0.3), Inches(0.65), BLUE)
        add_textbox(slide, stats_text,
                    ox + Inches(0.2), Inches(2.07),
                    half_w - Inches(0.4), Inches(0.6),
                    font_size=14, color=WHITE, wrap=True, font_category="body")

        # Divider
        add_rect(slide, ox + Inches(0.15), Inches(2.70),
                 half_w - Inches(0.3), Inches(0.025), GOLD)

        # Insight label
        add_textbox(slide, "CONTENT STRATEGY INSIGHT",
                    ox + Inches(0.15), Inches(2.80),
                    half_w - Inches(0.3), Inches(0.3),
                    font_size=10, bold=True, color=GOLD, font_category="caption")

        # Insight text — full, no truncation
        add_textbox(slide, insight or "Insufficient data for this company.",
                    ox + Inches(0.15), Inches(3.18),
                    half_w - Inches(0.3), Inches(3.72),
                    font_size=14, color=WHITE, wrap=True, font_category="body", max_words=110)


def slide_05x_content_insights(prs, all_data, analysis, your_company):
    """
    Generates 1 + ceil((n-1)/2) insight slides after slide 05:
      - Slide A: Full-width for YOUR company
      - Slide B: First two competitors side by side
      - Slide C: Next two competitors side by side (if any)
      ...etc
    """
    themes = analysis.get("content_themes", {})
    found  = [d for d in all_data if d["channel_found"]]

    your_data   = next((d for d in found if d["company_name"] == your_company), None)
    competitors = [d for d in found if d["company_name"] != your_company]

    # Slide A — your company full width
    if your_data:
        insight = themes.get(your_data["company_name"], "No insight data available.")
        _slide_content_insight_your_company(prs, your_data, insight)

    # Slides B, C, ... — competitors in pairs
    for i in range(0, len(competitors), 2):
        comp_a = competitors[i]
        comp_b = competitors[i + 1] if i + 1 < len(competitors) else None
        ins_a  = themes.get(comp_a["company_name"], "No insight data available.")
        ins_b  = themes.get(comp_b["company_name"], "No insight data available.") if comp_b else None
        _slide_content_insight_pair(prs, comp_a, ins_a, comp_b, ins_b)


def slide_06_posting_frequency(prs, all_data, insights):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_header(slide, "Posting Frequency & Consistency",
                     "Who is most active and on what cadence")

    found      = [d for d in all_data if d["channel_found"]]
    bar_max_w  = Inches(9.5)
    bar_h      = Inches(0.65)
    bar_gap    = Inches(0.85)
    bar_area_y = Inches(1.4)

    freqs    = [d.get("upload_frequency_days") or 999 for d in found]
    min_freq = min(freqs) if freqs else 1
    max_freq = max(freqs) if freqs else 999

    for i, d in enumerate(found):
        y    = bar_area_y + i * bar_gap
        freq = d.get("upload_frequency_days") or 999

        if max_freq == min_freq:
            bar_len = bar_max_w * 0.7
        else:
            bar_len = bar_max_w * (1 - (freq - min_freq) / (max_freq - min_freq)) * 0.85 + bar_max_w * 0.15

        add_textbox(slide, d["company_name"], Inches(0.4), y, Inches(2.2), bar_h,
                    font_size=14, bold=True, color=NAVY, font_category="body")
        add_rect(slide, Inches(2.8), y + Inches(0.1), bar_max_w, bar_h - Inches(0.2), LIGHTBLUE)
        bar_color = GREEN if i == 0 else BLUE
        add_rect(slide, Inches(2.8), y + Inches(0.1), bar_len, bar_h - Inches(0.2), bar_color)
        freq_label = f"Every {freq} days" if freq != 999 else "Unknown"
        label_x = Inches(2.8) + bar_len + Inches(0.1)
        label_x = min(label_x, Inches(11.0))  # Never overflow slide right edge
        add_textbox(slide, freq_label,
                    label_x, y + Inches(0.1), Inches(2.2), bar_h,
                    font_size=11, color=DARKGREY, font_category="caption")

    posting_insight = insights.get("posting_insight", "")
    add_rect(slide, Inches(0.4), Inches(5.8), Inches(12.5), Inches(1.4), NAVY)
    add_textbox(slide, "📊 " + posting_insight,
                Inches(0.6), Inches(5.9), Inches(12), Inches(1.2),
                font_size=14, color=WHITE, wrap=True, italic=True, font_category="body")


def slide_07_engagement(prs, all_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_header(slide, "Engagement Analysis",
                     "Average views, likes, and comments per video")

    found   = [d for d in all_data if d["channel_found"]]
    card_w  = Inches(12.5) / max(len(found), 1)
    start_x = Inches(0.35)

    for i, d in enumerate(found):
        x         = start_x + i * card_w
        eng       = d.get("engagement_rate", 0)
        eng_color = GREEN if eng > 2 else (GOLD if eng > 0.5 else RED)

        add_rect(slide, x, Inches(1.25), card_w - Inches(0.1), Inches(0.45), NAVY)
        add_textbox(slide, d["company_name"], x + Inches(0.05), Inches(1.3),
                    card_w - Inches(0.15), Inches(0.35),
                    font_size=11, bold=True, color=WHITE, font_category="caption")

        metrics = [
            ("Avg Views",       fmt_num(d["avg_views_per_video"]),    "per video"),
            ("Avg Likes",       fmt_num(d["avg_likes_per_video"]),    "per video"),
            ("Avg Comments",    fmt_num(d["avg_comments_per_video"]), "per video"),
            ("Engagement Rate", f"{eng}%",                            "likes+comments/views"),
        ]
        for j, (label, val, sub) in enumerate(metrics):
            my = Inches(1.8) + j * Inches(1.3)
            bg = BLUE if j < 3 else NAVY
            add_rect(slide, x + Inches(0.05), my, card_w - Inches(0.15), Inches(1.2), bg)
            add_textbox(slide, label,
                        x + Inches(0.1), my + Inches(0.05), card_w - Inches(0.2), Inches(0.3),
                        font_size=11, color=LIGHTBLUE, align=PP_ALIGN.CENTER, font_category="caption")
            add_textbox(slide, val,
                        x + Inches(0.1), my + Inches(0.35), card_w - Inches(0.2), Inches(0.45),
                        font_size=20, bold=True,
                        color=eng_color if j == 3 else WHITE,
                        align=PP_ALIGN.CENTER, font_category="title")
            add_textbox(slide, sub,
                        x + Inches(0.1), my + Inches(0.8), card_w - Inches(0.2), Inches(0.3),
                        font_size=11, color=LIGHTBLUE, align=PP_ALIGN.CENTER, font_category="caption")


def slide_08_gap_analysis(prs, analysis, insights):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_header(slide, "Content Gap Analysis",
                     "Topics competitors are covering that you are missing")

    gaps     = analysis.get("gaps", [])
    gap_text = insights.get("gap_analysis", "")

    add_rect(slide, Inches(0.4), Inches(1.3), Inches(12.5), Inches(1.3), WHITE)
    add_textbox(slide, gap_text, Inches(0.6), Inches(1.35), Inches(12), Inches(1.2),
                font_size=14, color=DARKGREY, wrap=True, font_category="body")

    add_textbox(slide, "Content Topics You Are Missing:", Inches(0.4), Inches(2.8),
                Inches(8), Inches(0.4),
                font_size=14, bold=True, color=NAVY, font_category="body")

    pill_x     = Inches(0.4)
    pill_y     = Inches(3.25)
    pill_h     = Inches(0.55)
    pill_gap   = Inches(0.1)
    max_row_w  = Inches(12.5)

    for gap in gaps:
        pill_w = max(Inches(1.5), Inches(len(gap) * 0.13 + 0.4))
        if pill_x + pill_w > Inches(0.4) + max_row_w:
            pill_x  = Inches(0.4)
            pill_y += pill_h + pill_gap
        add_rect(slide, pill_x, pill_y, pill_w, pill_h, RED)
        add_textbox(slide, gap.upper(),
                    pill_x + Inches(0.1), pill_y + Inches(0.1),
                    pill_w - Inches(0.2), Inches(0.35),
                    font_size=11, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, font_category="caption")
        pill_x += pill_w + pill_gap

    add_rect(slide, Inches(0.4), Inches(5.75), Inches(12.5), Inches(1.5), NAVY)
    add_textbox(slide, "💡 OPPORTUNITY", Inches(0.6), Inches(5.85), Inches(4), Inches(0.4),
                font_size=14, bold=True, color=GOLD, font_category="body")
    opp_text = (
        f"These {len(gaps)} content areas represent immediate opportunities. "
        "Being first to own these topics on YouTube establishes thought leadership "
        "and captures search traffic your competitors are ignoring."
    )
    add_textbox(slide, opp_text, Inches(0.6), Inches(6.3), Inches(12), Inches(0.8),
                font_size=14, color=WHITE, wrap=True, font_category="body")


def slide_09_recommendations(prs, insights, your_company):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_header(slide, "Video Marketing Recommendations",
                     f"Specific, data-backed actions for {your_company}")

    recs_text = insights.get("recommendations", "")
    recs      = [r.strip() for r in recs_text.split("\n") if r.strip() and r.strip()[0].isdigit()]
    icons     = ["🎯", "📈", "🎬", "🤝"]
    colors    = [BLUE, GREEN, GOLD, NAVY]
    labels    = ["Priority Action", "Growth Lever", "Content Strategy", "Engagement Tactic"]

    for i, rec in enumerate(recs[:4]):
        clean = rec.lstrip("0123456789. ").strip()
        x     = Inches(0.35) + (i % 2) * Inches(6.3)
        y     = Inches(1.3) + (i // 2) * Inches(2.9)

        add_rect(slide, x, y, Inches(6.1), Inches(2.7), colors[i] if i < len(colors) else BLUE)
        add_textbox(slide, icons[i] if i < len(icons) else "▶",
                    x + Inches(0.15), y + Inches(0.12), Inches(0.55), Inches(0.5),
                    font_size=20, color=WHITE, font_category="title")
        add_textbox(slide, labels[i] if i < len(labels) else f"Action {i+1}",
                    x + Inches(0.8), y + Inches(0.15), Inches(5.0), Inches(0.4),
                    font_size=16, bold=True, color=WHITE, font_category="body")
        add_textbox(slide, clean,
                    x + Inches(0.15), y + Inches(0.68), Inches(5.8), Inches(1.9),
                    font_size=14, color=WHITE, wrap=True, font_category="body")


def slide_10_scorecard(prs, analysis, all_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_header(slide, "Company Score Card",
                     "Composite ranking across all key video marketing metrics")

    ranked = analysis.get("ranked", [])
    dims   = ["Company", "Subscribers", "Avg Views", "Engagement", "Frequency", "TOTAL SCORE"]
    dim_w  = [Inches(2.8), Inches(1.7), Inches(1.7), Inches(1.9), Inches(1.8), Inches(2.0)]
    start_x = Inches(0.35)
    row_h   = Inches(0.7)
    y       = Inches(1.3)

    x = start_x
    for dim, dw in zip(dims, dim_w):
        add_rect(slide, x, y, dw - Inches(0.05), row_h, NAVY)
        add_textbox(slide, dim, x + Inches(0.05), y + Inches(0.15),
                    dw - Inches(0.1), Inches(0.45),
                    font_size=11, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, font_category="caption")
        x += dw

    for row_i, r in enumerate(ranked):
        y  = Inches(1.3) + row_h + row_i * row_h
        bg = LIGHTBLUE if row_i == 0 else (GREY if row_i % 2 == 0 else WHITE)
        x  = start_x
        scores = [
            r["company"],
            str(r["subscribers_score"]),
            str(r["views_score"]),
            str(r["engagement_score"]),
            str(r["frequency_score"]),
            f"⭐ {r['total_score']}/10"
        ]
        for si, (val, dw) in enumerate(zip(scores, dim_w)):
            add_rect(slide, x, y, dw - Inches(0.05), row_h, bg)
            fc = NAVY if si == 0 else (BLUE if si == 5 else DARKGREY)
            fb = si in (0, 5)
            add_textbox(slide, val, x + Inches(0.05), y + Inches(0.15),
                        dw - Inches(0.1), Inches(0.45),
                        font_size=11, bold=fb, color=fc,
                        align=PP_ALIGN.CENTER, font_category="caption")
            x += dw

    add_textbox(slide,
                "Scores normalised 1–10. Weighted: Views 30% · Subscribers 25% · Engagement 25% · Frequency 20%",
                Inches(0.35), Inches(6.9), Inches(12.5), Inches(0.4),
                font_size=11, color=DARKGREY, italic=True, font_category="caption")


def slide_11_summary(prs, your_company, insights, analysis):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, Inches(0), Inches(0), W, H, NAVY)
    add_rect(slide, Inches(0), Inches(0), Inches(0.5), H, BLUE)
    add_rect(slide, Inches(0), Inches(7.1), W, Inches(0.4), GOLD)

    add_textbox(slide, "Key Takeaways", Inches(1), Inches(0.6), Inches(10), Inches(0.7),
                font_size=32, bold=True, color=WHITE, font_category="title")
    add_rect(slide, Inches(1), Inches(1.35), Inches(8), Inches(0.04), GOLD)

    leader = analysis.get("leader", "N/A")
    gaps   = analysis.get("gaps", [])[:5]

    takeaways = [
        f"🏆 {leader} leads the competitive landscape in YouTube video marketing.",
        f"📊 Content gap identified — {len(analysis.get('gaps', []))} topics your competitors cover that you don't.",
        "🎬 Consistent upload frequency is the single biggest differentiator between leaders and laggards.",
        "💡 Engagement quality matters — high likes+comments signal content that builds audience trust.",
    ]
    if gaps:
        takeaways.append(f"🎯 Immediate opportunity: Create content on — {', '.join(gaps[:4])}.")

    for i, t in enumerate(takeaways):
        y = Inches(1.55) + i * Inches(0.72)
        add_textbox(slide, t, Inches(1), y, Inches(11), Inches(0.68),
                    font_size=16, color=WHITE, wrap=True, font_category="body")

    score_just = insights.get("score_justification", "")
    if score_just:
        add_rect(slide, Inches(1), Inches(5.45), Inches(11.5), Inches(1.55), BLUE)
        add_textbox(slide, score_just, Inches(1.2), Inches(5.52), Inches(11), Inches(1.4),
                    font_size=18, color=WHITE, wrap=True, italic=True, font_category="body")

    add_textbox(slide,
                f"Report generated for {your_company} · YouTube Competitor Intelligence Tool",
                Inches(1), Inches(7.1), Inches(11), Inches(0.35),
                font_size=11, color=DARKGREY, italic=True, font_category="caption")


# ---------------------------------------------------------------------------
# Master builder
# ---------------------------------------------------------------------------

def build_pptx(your_company, competitors, all_data, analysis, insights):
    """
    Builds the complete 11-slide PPTX and returns it as an in-memory BytesIO buffer.
    Called from app.py after the quality gate.
    """
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    report_date = datetime.now().strftime("%d %B %Y")

    slide_01_cover(prs, your_company, competitors, report_date)
    slide_02_executive_summary(prs, your_company, insights, analysis)
    slide_03_channel_overview(prs, all_data)
    slide_04_content_performance(prs, all_data)
    slide_05_content_topics(prs, all_data, analysis)
    slide_05x_content_insights(prs, all_data, analysis, your_company)  # 1 + N/2 insight slides
    slide_06_posting_frequency(prs, all_data, insights)
    slide_07_engagement(prs, all_data)
    slide_08_gap_analysis(prs, analysis, insights)
    slide_09_recommendations(prs, insights, your_company)
    slide_10_scorecard(prs, analysis, all_data)
    slide_11_summary(prs, your_company, insights, analysis)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
